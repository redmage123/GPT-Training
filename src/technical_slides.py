import collections
import collections.abc
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches,Pt
import sys

class GPT3Presentation:
    def __init__(self):
        self.prs = Presentation()

    def add_slide(self, title:str, bullet_text:list):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        bullet_shape = slide.shapes.placeholders[1]
        tf = bullet_shape.text_frame
        for text in bullet_text:
            p = tf.add_paragraph()
            p.text = text
            p.font.type = "Calibri"
            p.font.size = Pt(22)

    def generate_presentation(self,text_corpus:list): 
        for slide_dict in text_corpus:
            for key,value in slide_dict.items():
                self.add_slide(key,value)
        print ("Saving presentation")
        self.prs.save('GPT3-technical.pptx')
#text_frame = add_text_box(slide, text, Inches(0.5), Inches(6), Inches(9), Inches(2))

'''
def add_text_box(slide, text, left, top, width, height):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.text = text
    text_frame.paragraphs[0].font.name = "Calibri"
    text_frame.paragraphs[0].font.size = Pt(22)
    return text_frame
'''


prs = GPT3Presentation()
slides_text =  [{"Introduction:":["In today's world, technology has enabled us to interact with machines in a way that mimics human interaction.","This is made possible through Natural Language Processing (NLP) which involves the use of computer algorithms to analyze, understand, and generate human language.","In this presentation, we will discuss the various models used for NLP, including Recurrent Neural Networks (RNNs), Long Short-Term Memory Neural Networks (LSTMs), and Transformer Neural Networks.","We will also discuss the strengths and weaknesses of these models, as well as the use cases for Generative Pre-trained Transformer (GPT) models, including ChatGPT.","Finally, we will discuss the future directions of GPTs."]},
               {"What is Natural Language Processing?":["NLP is the branch of artificial intelligence that focuses on the interaction between computers and humans using natural language.","It involves the use of computer algorithms to process, analyze, and understand human language, including speech and text."]},
               {"Models used for NLP":["- Recurrent Neural Networks (RNNs)","- Long Short-Term Memory Neural Networks (LSTMs)","- Transformer Neural Networks"]},
               {"What are Recurrent Neural Networks?":["RNNs are a type of neural network that are designed to process sequential data.","They are commonly used for tasks such as speech recognition and machine translation.","In RNNs, the same function is applied to each element in a sequence and the output is passed from one step to the next."]},
               {"What are Long Short-Term Memory Neural Networks?":["LSTMs are a type of RNN that are designed to address the problem of vanishing gradients in RNNs.","They are used for tasks such as sentiment analysis and language generation.","In LSTMs, the network has a memory cell that can store information and use it to control the flow of information in the network."]},
               {"Weaknesses and Drawbacks of RNNs and LSTMs":["RNNs and LSTMs can be difficult to train and may not perform well when processing long sequences of data.","Additionally, they are not well suited for parallel processing, which can lead to slow performance."]},
               {"What is a Transformer Neural Network?":["Transformer Neural Networks are a type of neural network that is designed for NLP tasks, such as machine translation and question answering.","They are based on the concept of attention and use self-attention to analyze the relationships between different elements in a sequence."]},
               {"How Transformers Work":["Transformers work by using self-attention to analyze the relationships between different elements in a sequence.","They use multiple attention heads to analyze the relationships between different elements and to make predictions based on the information they have seen."]},
               {"'Attention is All You Need' Paper":["The paper 'Attention is All You Need' introduced the Transformer Neural Network model and demonstrated its effectiveness for NLP tasks.","It showed that the Transformer model was able to outperform traditional RNN and LSTM models on various NLP tasks."]},
               {"Attention and Self-Attention in Transformers":["Attention and self-attention are key concepts in the Transformer model.","Attention allows the network to focus on specific parts of the input and use that information to make predictions.","Self-attention allows the network to analyze the relationships between different elements in a sequence and make predictions based on that information."]},
               {"Use Cases for GPTs":["GPTs are used for a variety of NLP tasks, including machine translation, question answering, and language generation.","They are also used in chatbots and conversational AI, such as ChatGPT."]},
               {"What is ChatGPT?":["ChatGPT is a conversational AI model developed by OpenAI that uses the GPT architecture.","It is designed to respond to human input in a conversational manner and can be used for tasks such as customer service, virtual assistants, and chatbots."]},
               {"How ChatGPT Works":["ChatGPT works by using the GPT architecture to generate a response based on the input it receives.","The model has been pre-trained on a large corpus of text data, which allows it to understand the context of a conversation and generate appropriate responses."]},
               {"Future Directions for GPTs":["The field of NLP is rapidly evolving and there is a growing interest in the use of GPTs for various applications.","In the future, we can expect to see GPTs being used for even more complex NLP tasks, such as sentiment analysis and natural language generation.","Additionally, there is a growing interest in the use of GPTs for tasks such as recommendation systems and content creation."]},
               {"Conclusion":["In conclusion, GPTs and ChatGPT are powerful models that are changing the way we interact with computers using natural language.","They have been trained on large amounts of data, which allows them to understand the context of a conversation and generate appropriate responses.","With the growing interest in NLP and the increasing use of GPTs for various applications, we can expect to see continued advancements in this field in the years to come."]}]
prs.generate_presentation(slides_text)               