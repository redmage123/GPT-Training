import yaml
import collections
import collections.abc
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches,Pt
import sys
import os
from os import path

class GPT3Presentation:
    def __init__(self, yaml_file: str,output_file:str):
        self.prs = Presentation()
        self.slides_text = []
        self.yaml_file = yaml_file
        self.content = None
        self.font_type = "Calibri"
        self.font_size = Pt(22)
        self.slides_text = []
        self.output_file = output_file
        
            
    def load_yaml_file(self) :
        self.yaml_file = "/home/bbrelin/src/repos/GPT-Training/src/sample_slide.yaml"
        with open (self.yaml_file ,"r") as yf:
            try:
                self.content = yaml.safe_load(yf)
            except yaml.YAMLError as exc:  
                print (exc)
                sys.exit(1)
        #print (self.content)
        #sys.exit(0)
        return
    
    def parse_yaml_file(self):
        count = 0
        if self.content == None:
            assert("No data in yaml file!")
        for slide in self.content['slides']:
            if 'font' in slide:
                self.font = slide['font']
                continue
            elif 'fontsize' in slide:
                self.fontsize = slide['fontsize']
                continue
            else:
                d = {'title':slide['title'],'bullet_points':slide['bullet_points']}
#                d = {slide['title']:slide['bullet_points']}
                self.slides_text.append(d)
        return
        

    @staticmethod
    def extract_dicts(data):
        data = list(data)
        return [{data[i][0]: data[i][1], data[i + 1][0]: data[i + 1][1]} for i in range(0, len(data), 2)]
    
    def add_slide(self,ys):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title

        #print (ys.items())
        #sys.exit(1)
        for slide_dict  in self.extract_dicts(ys.items()):
            title_shape.text = slide_dict['title']
            bullet_shape = slide.shapes.placeholders[1]
            tf = bullet_shape.text_frame
            for text in slide_dict['bullet_points']:
                p = tf.add_paragraph()
                p.text = text
                p.font.type = self.font
                p.font.size = self.font_size

    def generate_presentation(self): 
        self.load_yaml_file()
        self.parse_yaml_file() 
        for ys in self.slides_text:
            self.add_slide(ys)
        print ("Saving presentation")
        self.prs.save('GPT3-technical.pptx')

class ReadConfig:
    def __init__(self, config_file_path:str):
        self.config_file_path = config_file_path
        
    def read_config(self) -> dict:
        config_dict = {}
        try: 
            with open(self.config_file_path,'r') as cf:
                for line in cf:
                    line_list = line.replace(' ','').split('=')
                    config_dict[line_list[0]] = line_list[1]
        except FileNotFoundError:
            raise FileNotFoundError
        return config_dict
if __name__ == '__main__':

    try: 
        cf = os.getcwd() + "/src/" + "Config.py"
        rf = ReadConfig(cf) 
        config_dict = rf.read_config()
    except FileNotFoundError:
        sys.stderr.write("Error:  No config file found!")
        sys.exit(1)

    prs = GPT3Presentation(config_dict['yaml_file'],config_dict['output_file'])
    prs.generate_presentation()               