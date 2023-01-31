from pptx import Presentation
from pptx.util import Inches

class GPT3Presentation:
    def __init__(self):
        self.prs = Presentation()
        self.title = ""
        self.font_size = 24

    def add_slide(self, title, bullet_text):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        bullet_shape = slide.shapes.placeholders[1]
        tf = bullet_shape.text_frame
        for text in bullet_text:
            tf.add_paragraph().text = text

    def generate_presentation(self):
        self.add_slide("Introduction to GPT and ChatGPT", ["What is GPT-3", "What is ChatGPT"])
        self.add_slide("Non-technical introduction to Natural Language Processing", ["What is NLP", "NLP applications"])
        self.add_slide("History of Natural Language Processing", ["Early NLP systems", "Recent developments"])
        self.add_slide("Non-technical introduction to Neural Networks", ["What are neural networks", "Applications of neural networks"])
        self.add_slide("History of Neural Networks", ["Perceptron", "Backpropagation"])
        self.add_slide("Non-technical introduction to Transformer Neural Networks", ["What are transformer networks", "Attention is all you need"])
        self.add_slide("A history of the OpenAI Foundation", ["OpenAI's mission", "Notable achievements"])
        self.add_slide("Use cases for GPT-3 and ChatGPT", ["Natural language generation", "Chatbots"])
        self.add_slide("Drawbacks and challenges in using ChatGPT", ["Limitations of GPT-3", "Ethical concerns"])
        self.add_slide("What is the future of GPT?", ["Potential developments", "Potential impact"])
        self.add_slide("Summary and Conclusions", ["Key points", "Future directions"])
        self.prs.save('GPT3Presentation.pptx')
